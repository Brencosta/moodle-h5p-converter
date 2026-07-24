
#!/usr/bin/env python3
"""
h5p_to_pptx.py  (versão robusta)
 
Converte um arquivo .h5p do tipo "Course Presentation" (H5P.CoursePresentation)
para um arquivo .pptx editável, preservando:
  - fundo do slide (cor ou imagem)
  - texto com negrito, itálico, cor, alinhamento e listas (<ul>/<ol>)
  - caixas de texto com cor de fundo (ex: faixas de título coloridas)
  - imagens
 
Elementos interativos (quiz, vídeo, drag-and-drop, botões de navegação etc.)
não têm equivalente em slide estático: são ignorados silenciosamente e
listados num resumo no final, para você saber o que ficou de fora.
 
Uso:
    python h5p_to_pptx.py caminho/para/arquivo.h5p [saida.pptx]
 
Requisitos:
    pip install python-pptx beautifulsoup4
"""
 
import sys
import os
import re
import json
import zipfile
import shutil
import tempfile
from collections import Counter
 
from bs4 import BeautifulSoup, NavigableString, Tag
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
 
SLIDE_WIDTH_EMU = Emu(9144000)    # 10"
SLIDE_HEIGHT_EMU = Emu(5143500)   # 16:9
 
# Bibliotecas que sabemos renderizar de verdade
LIBS_TEXTO = ("H5P.AdvancedText", "H5P.Text")
LIBS_IMAGEM = ("H5P.Image",)
 
# Bibliotecas que são só navegação/interação e podem ser ignoradas sem aviso
LIBS_IGNORAR_SILENCIOSO = ("H5P.ContinuousText",)  # ajuste se aparecer algo assim
 
ignorados_relatorio = Counter()
 
 
# --------------------------------------------------------------------------
# Utilidades
# --------------------------------------------------------------------------
 
def extrair_h5p(caminho_h5p, pasta_destino):
    with zipfile.ZipFile(caminho_h5p, "r") as z:
        z.extractall(pasta_destino)
 
 
def carregar_content_json(pasta_extraida):
    caminho = os.path.join(pasta_extraida, "content", "content.json")
    if not os.path.exists(caminho):
        raise FileNotFoundError(
            "content/content.json não encontrado — confira se o .h5p é "
            "realmente do tipo Course Presentation."
        )
    with open(caminho, "r", encoding="utf-8") as f:
        return json.load(f)
 
 
def pct_x(valor_percentual):
    return Emu(int(SLIDE_WIDTH_EMU * (valor_percentual / 100.0)))
 
 
def pct_y(valor_percentual):
    return Emu(int(SLIDE_HEIGHT_EMU * (valor_percentual / 100.0)))
 
 
def parse_cor(valor):
    """Converte '#rrggbb' ou 'rgb(r,g,b)' em RGBColor. Retorna None se falhar."""
    if not valor:
        return None
    valor = valor.strip()
    m = re.match(r"#([0-9a-fA-F]{6})$", valor)
    if m:
        return RGBColor.from_string(m.group(1))
    m = re.match(r"#([0-9a-fA-F]{3})$", valor)
    if m:
        h = "".join(c * 2 for c in m.group(1))
        return RGBColor.from_string(h)
    m = re.match(r"rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)", valor)
    if m:
        r, g, b = (int(x) for x in m.groups())
        return RGBColor(r, g, b)
    return None
 
 
def parse_style_attr(style_str):
    """Transforma 'color: red; font-weight: bold' num dict simples."""
    resultado = {}
    if not style_str:
        return resultado
    for parte in style_str.split(";"):
        if ":" not in parte:
            continue
        chave, valor = parte.split(":", 1)
        resultado[chave.strip().lower()] = valor.strip()
    return resultado
 
 
def achar_arquivo_imagem(pasta_extraida, caminho_relativo):
    """Tenta localizar o arquivo de imagem de forma tolerante a variações de path."""
    candidatos = [
        os.path.join(pasta_extraida, "content", caminho_relativo),
        os.path.join(pasta_extraida, "content", caminho_relativo.lstrip("/")),
        os.path.join(pasta_extraida, caminho_relativo),
    ]
    for c in candidatos:
        if os.path.exists(c):
            return c
    return None
 
 
# --------------------------------------------------------------------------
# Parser de HTML rico -> lista de parágrafos com runs formatados
# --------------------------------------------------------------------------
 
TAMANHO_PADRAO = 18
TAMANHOS_TITULO = {"h1": 32, "h2": 26, "h3": 22, "h4": 20, "h5": 18, "h6": 16}
 
 
def _extrair_runs(node, estilo_atual, runs_saida):
    """Percorre recursivamente os filhos, acumulando texto formatado em runs_saida."""
    if isinstance(node, NavigableString):
        texto = str(node)
        if texto.strip("\u00a0 \t\n\r") != "" or texto == "\u00a0":
            runs_saida.append((texto, dict(estilo_atual)))
        return
 
    if not isinstance(node, Tag):
        return
 
    novo_estilo = dict(estilo_atual)
    tag = node.name.lower()
 
    if tag in ("strong", "b"):
        novo_estilo["bold"] = True
    elif tag in ("em", "i"):
        novo_estilo["italic"] = True
    elif tag == "u":
        novo_estilo["underline"] = True
    elif tag in TAMANHOS_TITULO:
        novo_estilo["bold"] = True
        novo_estilo["size"] = TAMANHOS_TITULO[tag]
 
    css = parse_style_attr(node.get("style", ""))
    if "color" in css:
        cor = parse_cor(css["color"])
        if cor:
            novo_estilo["color"] = cor
    if css.get("font-weight") in ("bold", "700", "600"):
        novo_estilo["bold"] = True
    if css.get("font-style") == "italic":
        novo_estilo["italic"] = True
    if "text-decoration" in css and "underline" in css["text-decoration"]:
        novo_estilo["underline"] = True
 
    for filho in node.children:
        _extrair_runs(filho, novo_estilo, runs_saida)
 
 
def html_para_paragrafos(html):
    """
    Converte HTML do H5P (AdvancedText) numa lista de parágrafos:
    [{"runs": [(texto, estilo), ...], "align": "left"/"center"/"right", "bullet": "•"/"1."/None}]
    """
    if not html or not html.strip():
        return []
 
    soup = BeautifulSoup(html, "html.parser")
    paragrafos = []
 
    blocos = soup.find_all(["p", "h1", "h2", "h3", "h4", "h5", "h6", "li", "div"], recursive=True)
    # Evita pegar blocos aninhados duas vezes (ex: <li><p>texto</p></li>)
    blocos_top = [b for b in blocos if not any(b in outro.find_all() for outro in blocos if outro is not b)]
 
    if not blocos_top:
        # HTML sem tags de bloco reconhecíveis — trata tudo como um parágrafo só
        runs = []
        _extrair_runs(soup, {}, runs)
        if runs:
            paragrafos.append({"runs": runs, "align": None, "bullet": None})
        return paragrafos
 
    for bloco in blocos_top:
        runs = []
        _extrair_runs(bloco, {}, runs)
        if not runs:
            continue
 
        align = None
        css = parse_style_attr(bloco.get("style", ""))
        if css.get("text-align") in ("left", "center", "right", "justify"):
            align = css["text-align"]
 
        bullet = None
        if bloco.name == "li":
            pai = bloco.find_parent(["ul", "ol"])
            bullet = "1." if (pai and pai.name == "ol") else "•"
 
        paragrafos.append({"runs": runs, "align": align, "bullet": bullet})
 
    return paragrafos
 
 
ALINHAMENTO_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}
 
 
def preencher_textbox(caixa, paragrafos):
    tf = caixa.text_frame
    tf.word_wrap = True
    primeiro = True
 
    for par in paragrafos:
        p = tf.paragraphs[0] if primeiro else tf.add_paragraph()
        primeiro = False
 
        if par["align"]:
            p.alignment = ALINHAMENTO_MAP.get(par["align"])
 
        prefixo = f"{par['bullet']} " if par["bullet"] else ""
        primeiro_run = True
 
        for texto, estilo in par["runs"]:
            if primeiro_run and prefixo:
                texto = prefixo + texto
            primeiro_run = False
 
            texto = texto.replace("\u00a0", " ")
            if texto == "":
                continue
 
            run = p.add_run()
            run.text = texto
            run.font.size = Pt(estilo.get("size", TAMANHO_PADRAO))
            run.font.bold = estilo.get("bold", False)
            run.font.italic = estilo.get("italic", False)
            run.font.underline = estilo.get("underline", False)
            if "color" in estilo:
                run.font.color.rgb = estilo["color"]
 
 
# --------------------------------------------------------------------------
# Fundo do slide
# --------------------------------------------------------------------------
 
def aplicar_fundo_slide(slide, bg_info, pasta_extraida):
    if not bg_info:
        return
 
    fill_type = bg_info.get("fillType")
 
    if fill_type == "color" or ("colorSettings" in bg_info and not fill_type):
        cor_hex = (bg_info.get("colorSettings") or {}).get("color")
        cor = parse_cor(cor_hex)
        if cor:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = cor
            return
 
    if fill_type == "image" or ("imageSettings" in bg_info and not fill_type):
        img = (bg_info.get("imageSettings") or {}).get("fill", {})
        caminho_rel = img.get("path")
        if caminho_rel:
            caminho_abs = achar_arquivo_imagem(pasta_extraida, caminho_rel)
            if caminho_abs:
                slide.shapes.add_picture(
                    caminho_abs, 0, 0, width=SLIDE_WIDTH_EMU, height=SLIDE_HEIGHT_EMU
                )
                # manda a imagem de fundo pro fundo da pilha de elementos
                spTree = slide.shapes._spTree
                pic = slide.shapes[-1]._element
                spTree.remove(pic)
                spTree.insert(2, pic)
 
 
# --------------------------------------------------------------------------
# Elementos
# --------------------------------------------------------------------------
 
def cor_de_fundo_do_html(html):
    """Se o HTML tem um wrapper com background-color inline, extrai a cor."""
    if not html:
        return None
    soup = BeautifulSoup(html, "html.parser")
    for tag in soup.find_all(True):
        css = parse_style_attr(tag.get("style", ""))
        if "background-color" in css or "background" in css:
            cor = parse_cor(css.get("background-color") or css.get("background"))
            if cor:
                return cor
    return None
 
 
def adicionar_texto(slide, elemento, left, top, width, height):
    params = elemento["action"].get("params", {})
    html = params.get("text", "")
    paragrafos = html_para_paragrafos(html)
    if not paragrafos:
        return
 
    cor_fundo = cor_de_fundo_do_html(html)
    opacidade = elemento.get("backgroundOpacity", 0)
 
    if cor_fundo or opacidade:
        caixa_fundo = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        caixa_fundo.line.fill.background()
        caixa_fundo.fill.solid()
        caixa_fundo.fill.fore_color.rgb = cor_fundo or RGBColor(0xFF, 0xFF, 0xFF)
        caixa_fundo.shadow.inherit = False
        tf = caixa_fundo.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(45720)
        tf.margin_right = Emu(45720)
        preencher_textbox(caixa_fundo, paragrafos)
    else:
        caixa = slide.shapes.add_textbox(left, top, width, height)
        preencher_textbox(caixa, paragrafos)
 
 
def adicionar_imagem(slide, elemento, left, top, width, height, pasta_extraida):
    params = elemento["action"].get("params", {})
    arquivo_info = params.get("file", {})
    caminho_relativo = arquivo_info.get("path", "") if isinstance(arquivo_info, dict) else ""
    if not caminho_relativo:
        ignorados_relatorio["H5P.Image (sem arquivo)"] += 1
        return
    caminho_imagem = achar_arquivo_imagem(pasta_extraida, caminho_relativo)
    if caminho_imagem:
        slide.shapes.add_picture(caminho_imagem, left, top, width, height)
    else:
        ignorados_relatorio["H5P.Image (arquivo não encontrado)"] += 1
 
 
def adicionar_elemento(slide, elemento, pasta_extraida):
    left = pct_x(elemento.get("x", 0))
    top = pct_y(elemento.get("y", 0))
    width = pct_x(elemento.get("width", 30))
    height = pct_y(elemento.get("height", 20))
 
    action = elemento.get("action")
    if not action:
        # elemento de navegação (goToSlide) ou botão puro — ignora silenciosamente
        if "goToSlide" in elemento:
            return
        ignorados_relatorio["elemento sem 'action' desconhecido"] += 1
        return
 
    library = action.get("library", "")
    nome_lib = library.split(" ")[0]
 
    try:
        if nome_lib in LIBS_TEXTO:
            adicionar_texto(slide, elemento, left, top, width, height)
        elif nome_lib in LIBS_IMAGEM:
            adicionar_imagem(slide, elemento, left, top, width, height, pasta_extraida)
        elif nome_lib in LIBS_IGNORAR_SILENCIOSO:
            return
        else:
            ignorados_relatorio[nome_lib] += 1
    except Exception as e:
        ignorados_relatorio[f"{nome_lib} (erro: {e})"] += 1
 
 
# --------------------------------------------------------------------------
# Conversão principal
# --------------------------------------------------------------------------
 
def converter(caminho_h5p, caminho_saida):
    pasta_temp = tempfile.mkdtemp(prefix="h5p_extract_")
    try:
        extrair_h5p(caminho_h5p, pasta_temp)
        dados = carregar_content_json(pasta_temp)
 
        slides_h5p = dados.get("presentation", {}).get("slides", [])
        if not slides_h5p:
            raise ValueError(
                "Nenhum slide encontrado em content.json — verifique se é "
                "mesmo um Course Presentation."
            )
 
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH_EMU
        prs.slide_height = SLIDE_HEIGHT_EMU
        layout_em_branco = prs.slide_layouts[6]
 
        for slide_h5p in slides_h5p:
            slide = prs.slides.add_slide(layout_em_branco)
            aplicar_fundo_slide(slide, slide_h5p.get("slideBackgroundSelector"), pasta_temp)
            for elemento in slide_h5p.get("elements", []):
                adicionar_elemento(slide, elemento, pasta_temp)
 
        prs.save(caminho_saida)
 
        print(f"OK: {len(slides_h5p)} slide(s) convertido(s) para: {caminho_saida}")
        if ignorados_relatorio:
            print("\nElementos NÃO renderizados (sem equivalente em slide estático):")
            for nome, qtd in ignorados_relatorio.most_common():
                print(f"  - {nome}: {qtd}x")
        else:
            print("Todos os elementos foram convertidos sem problemas.")
 
    finally:
        shutil.rmtree(pasta_temp, ignore_errors=True)
 
 
if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Uso: python h5p_to_pptx.py arquivo.h5p [saida.pptx]")
        sys.exit(1)
 
    entrada = sys.argv[1]
    saida = sys.argv[2] if len(sys.argv) > 2 else os.path.splitext(entrada)[0] + ".pptx"
 
    converter(entrada, saida)
 
